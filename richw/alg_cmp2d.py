#! /u/aurora-r0/richw/pkgs/miniforge3/envs/dist-s1/bin/python -i

from pptx import Presentation
from pptx.util import Inches
import os
import sys
import inspect
import numpy as np
import geopandas as gpd
import rasterio
from rasterio.crs import CRS
from rasterio.windows import Window
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from shapely.geometry import Point
import pandas as pd
from dem_stitcher.geojson_io import read_geojson_gzip
from dem_stitcher.geojson_io import to_geojson_gzip
from shapely.geometry import Polygon, shape
import asf_search as asf
import concurrent.futures
from tqdm import tqdm
import itertools
from datetime import datetime, timedelta
from dateutil.parser import parse
import ast

import distmetrics
import data_fcns
import plot_fcns
import test_setup
import alg_fcns

dir_base = '/home/richw/src/opera'
try:
  df_rtc,df_val_bursts_subset,df_sites_subset,df_slcs,df_burst,dir_rtc_site_data,dir_rtc_site_plots,basename_rtc_site_data = test_setup.setup(dir_base)
  dir_rtc_data = '/u/aurora-r0/cmarshak/dist-s1-research/marshak/6_torch_dataset/opera_rtc_data'
except:
  exc_type,exc_value,exc_traceback = sys.exc_info()
  tb = exc_traceback
  stack = inspect.trace()
  lcls = locals()
  slcls = tb.tb_next.tb_frame.f_locals
  raise Exception("Error in test_setup")


print('Iterating over site subset')

fstr = "%Y%m%dT%H%M%SZ"
fstr2 = "%Y-%m-%d %H:%M:%S"

delta_10min = datetime(2022,11,26,14,10,0) - datetime(2022,11,26,14,0,0)

prs = Presentation()
blank_slide_layout_a = prs.slide_layouts[6]
#prs_d = Presentation()
#blank_slide_layout_d = prs_d.slide_layouts[6]

fstr = "%Y%m%dT%H%M%SZ"
td_lookback = timedelta(days=365)
td_halfwindow = timedelta(days=12)

Nsubset = len(df_val_bursts_subset)

#for index, df_row in df_val_bursts_subset.iterrows():
#for index in range(Nsubset):
for index in range(0,1):
#  index,df_row = next(df_val_bursts_subset.iterrows())
  # Need burst id using uppercase and dashes
  burst_id = df_val_bursts_subset.loc[index,'jpl_burst_id']
  site_id = df_val_bursts_subset.loc[index,'site_id']
  # if this csv file does not exist, it's data file will also have not data
  csvname_rtc_site_data = (dir_rtc_site_data + '/' + basename_rtc_site_data +
    str(site_id) + '_burst_' + burst_id + '.csv')
  plotname_rtc_site_data = (dir_rtc_site_plots + '/' + basename_rtc_site_data +
    str(site_id) + '_burst_' + burst_id + '.png')
  df_site = df_sites_subset[df_sites_subset.site_id
    == site_id].reset_index(drop=True)
  lon = df_site.geometry.x[0]
  lat = df_site.geometry.y[0]
  # Site locations use long,lat defined on WGS84 ellipsoid which is the basis
  # of the EPSG:4326 coordinate reference system (CRS)
  gdf1 = gpd.GeoDataFrame(df_site.geometry,crs="EPSG:4326")
  # RTC data stored in directories named by burst id (uppercase,dashes)
  dir_burst1 = os.path.join(dir_rtc_data,burst_id)
  # RTC dataframe also uses burst id label (uppercase,dashes)
  # Note: order of entries of df_rtc_ts differs from order of filename_rtc
  df_rtc_ts = df_rtc[df_rtc.jpl_burst_id == burst_id].reset_index(drop=True)
  if len(df_rtc_ts) == 0:
    # This jpl_burst_id does not have data for this site_id
    print(f"id: {df_row['site_id']}, jpl_burst_id: {df_row['jpl_burst_id']} not present")
    continue

  print(f"site id: {site_id}, burst_id: {burst_id}")

  # Determine orbit direction
  df_burst_subset = df_burst.loc[df_burst['jpl_burst_id'] == burst_id]
  if len(df_burst_subset) == 1:
    odir = df_burst_subset['orbit_pass'].item()
  elif len(df_burst_subset) == 0:
    odir = 'NONE'
  else:
    odir = 'NONE'
    print('Warning: multiple burst id match')


  df_rtc_ts_wind = pd.read_csv(csvname_rtc_site_data).drop(columns=
    ['Unnamed: 0'])
  df_rtc_ts_wind['datetime'] = pd.to_datetime(df_rtc_ts_wind['datetime'])

  df_rtc_ts = df_rtc[df_rtc.jpl_burst_id == burst_id].reset_index(drop=True)
  df_rtc_ts['acq_datetime'] = pd.to_datetime(df_rtc_ts['acq_datetime'])

  # Convert srting to list
  df_rtc_ts_wind['vv_l'] = df_rtc_ts_wind['vv'].apply(data_fcns.convert_to_list)
  df_rtc_ts_wind['vh_l'] = df_rtc_ts_wind['vh'].apply(data_fcns.convert_to_list)
  df_rtc_ts_wind['vv/vh_l'] = df_rtc_ts_wind['vv/vh'].apply(data_fcns.convert_to_list)

  df_rtc_ts_wind['vv_avg'] = df_rtc_ts_wind['vv_l'].apply(lambda x: sum(x) / len(x) if isinstance(x, list) else x)#.rolling(3).sum()
  df_rtc_ts_wind['vh_avg'] = df_rtc_ts_wind['vh_l'].apply(lambda x: sum(x) / len(x) if isinstance(x, list) else x)#.rolling(3).sum()
  df_rtc_ts_wind['vh_center'] = df_rtc_ts_wind['vh_l'].apply(lambda x: np.array(x[0:9]).reshape((3, 3))[1, 1] if isinstance(x, list) else x)#.rolling(3).sum()
  df_rtc_ts_wind['vv/vh_avg'] = df_rtc_ts_wind['vv/vh_l'].apply(lambda x: sum(x) / len(x) if isinstance(x, list) else x)

  try:
    # Each file has data from one datetime.
    # Iterate over all and load those needed by the lookback window
    # and the post window.
    datetimes = []
    filelist = []
    for filename_rtc in os.listdir(dir_burst1):
      str_channel = filename_rtc[-6:-4]
      str_datetime = filename_rtc[32:48]
      datetime1 = datetime.strptime(str_datetime,fstr)
      datetimes.append(datetime1)
      filelist.append(filename_rtc)

    dist_ob_list = []
    dist_ob_dt = []
    pre_use = []
    for i,dt in enumerate(datetimes):
      dt_ref1 = dt - td_lookback - td_halfwindow
      dt_ref2 = dt - td_lookback + td_halfwindow
      iuse = data_fcns.indices(datetimes,
        lambda x: (x>=dt_ref1 and x <= dt_ref2))
      print(f"{dt}: {iuse}")
      pre_use.append(iuse)

    vv_data = []
    vh_data = []
    for i,filename_rtc in enumerate(filelist):
      iuse = pre_use[i]
      if len(iuse) > 0:
        for ii in iuse:
          filename_rtc1 = filelist[ii]
          file_path = os.path.join(dir_burst1, filename_rtc1)
          if os.path.isfile(file_path):
            str_channel = filename_rtc[-6:-4]
            row = -1
            col = -1
            with rasterio.open(file_path) as rtc_dataset:
              if row == -1:
                # burst coordinate system remains the same for all datetimes
                crs = rtc_dataset.crs
                # Reproject site location into current CRS
                gdf1_re = gdf1.to_crs(crs)
                x = gdf1_re.geometry.x[0]
                y = gdf1_re.geometry.y[0]
                row,col = rtc_dataset.index(x,y)
                print(f"channel: {str_channel} datetime: {str_datetime} r,c: {row},{col}")
    
              #band1 = rtc_dataset.read(1)
              # The window will actually load a chunk which is 512x512, and then
              # subselect the 3x3 matrix centered on row,col
              window = Window(col-1,row-1,3,3)
              band3x3 = rtc_dataset.read(1,window=window)
              if str_channel == "VV":
                vv_data.append(band3x3)
              else:
                vh_data.append(band3x3)

    #y_obs = alg_fcns.window_mahalanobis_1d_workflow(df_rtc_ts_wind['vh_l'])
    #dist_objs = alg_fcns.window_mahalanobis_1d_running(df_rtc_ts_wind['vh_l'])
    #dist_vals = [d_obj.dist[1, 1] if d_obj is not None else np.nan for d_obj in dist_objs]
    #dist_dts = [df_rtc_ts_wind['datetime'][k] for (k, d_obj) in enumerate(dist_objs) if d_obj is not None]
    #dist_plot = [dist_obj.dist[1, 1] for dist_obj in dist_objs if dist_obj is not None]

    #slide = prs.slides.add_slide(blank_slide_layout_a)

    #left = Inches(0.5)
    #top = Inches(1)
    ##height = Inches(6)
    #width = Inches(9)
    #pic = slide.shapes.add_picture('tmp.png',left,top,width=width)

  except Exception as e:
    exc_type,exc_value,exc_traceback = sys.exc_info()
    tb = exc_traceback
    stack = inspect.trace()
    lcls = locals()
    slcls = tb.tb_next.tb_frame.f_locals
    print(f"Alg Error in site id: {site_id}, burst_id: {burst_id}")
    print(exc_value)

#prs.save('rtc2d.pptx')
