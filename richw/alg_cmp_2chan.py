#! /u/aurora-r0/richw/pkgs/miniforge3/envs/dist-s1/bin/python -i

from pptx import Presentation
from pptx.util import Inches
import os
import sys
import inspect
import traceback
import numpy as np
import geopandas as gpd
import rasterio
from rasterio.crs import CRS
from rasterio.windows import Window
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from shapely.geometry import Point
import pandas as pd
from dem_stitcher.geojson_io import read_geojson_gzip
from dem_stitcher.geojson_io import to_geojson_gzip
from shapely.geometry import Polygon, shape
import asf_search as asf
import concurrent.futures
from tqdm import tqdm
import itertools
from datetime import datetime, timedelta
from dateutil.parser import parse
import ast

import distmetrics
import data_fcns
import plot_fcns
import test_setup
import alg_fcns

dir_base = '/home/richw/src/opera'
try:
  df_rtc,df_val_bursts_subset,df_sites_subset,df_slcs,df_burst,dir_rtc_site_data,dir_rtc_site_plots,basename_rtc_site_data = test_setup.setup(dir_base)
except:
  exc_type,exc_value,exc_traceback = sys.exc_info()
  tb = exc_traceback
  stack = inspect.trace()
  lcls = locals()
  slcls = tb.tb_next.tb_frame.f_locals
  raise Exception("Error in test_setup")


print('Iterating over site subset')

fstr = "%Y%m%dT%H%M%SZ"
fstr2 = "%Y-%m-%d %H:%M:%S"

delta_dt_all = []
delta_dt_all_maxabs = []
delta_vv_all = []
delta_vv_all_maxabs = []
delta_vh_all = []
delta_vh_all_maxabs = []
max_delta_dt_allsites = []
max_delta_vv_allsites = []
max_delta_vh_allsites = []
nan_count = []
lendat_mismatch_allsites = []
site_ids_compared = []
burst_ids_compared = []
delta_10min = datetime(2022,11,26,14,10,0) - datetime(2022,11,26,14,0,0)
 
N_WORKERS_IO = 10
n_pre_img = 5
n_post_imgs_to_confirm = 3
lookback_length_days = 365
threshold = 3.0

prs = Presentation()
blank_slide_layout_a = prs.slide_layouts[6]
#prs_d = Presentation()
#blank_slide_layout_d = prs_d.slide_layouts[6]

Nsubset = len(df_val_bursts_subset)

#for index, df_row in df_val_bursts_subset.iterrows():
for index in range(Nsubset):
#for index in range(0,1):
#  index,df_row = next(df_val_bursts_subset.iterrows())
  # Need burst id using uppercase and dashes
  burst_id = df_val_bursts_subset.loc[index,'jpl_burst_id']
  site_id = df_val_bursts_subset.loc[index,'site_id']
  # if this csv file does not exist, it's data file will also have not data
  csvname_rtc_site_data = (dir_rtc_site_data + '/' + basename_rtc_site_data +
    str(site_id) + '_burst_' + burst_id + '.csv')
  plotname_rtc_site_data = (dir_rtc_site_plots + '/' + basename_rtc_site_data +
    str(site_id) + '_burst_' + burst_id + '.png')
  df_site = df_sites_subset[df_sites_subset.site_id
    == site_id].reset_index(drop=True)
  lon = df_site.geometry.x[0]
  lat = df_site.geometry.y[0]
  # Site locations use long,lat defined on WGS84 ellipsoid which is the basis
  # of the EPSG:4326 coordinate reference system (CRS)
  gdf1 = gpd.GeoDataFrame(df_site.geometry,crs="EPSG:4326")
  # RTC dataframe also uses burst id label (uppercase,dashes)
  # Note: order of entries of df_rtc_ts differs from order of filename_rtc
  df_rtc_ts = df_rtc[df_rtc.jpl_burst_id == burst_id].reset_index(drop=True)
  df_rtc_ts['acq_datetime'] = pd.to_datetime(df_rtc_ts['acq_datetime'])
  if len(df_rtc_ts) == 0:
    # This jpl_burst_id does not have data for this site_id
    print(f"id: {df_row['site_id']}, jpl_burst_id: {df_row['jpl_burst_id']} not present")
    continue
  acq_dt_l = df_rtc_ts.acq_datetime.tolist()

  print(f"site id: {site_id}, burst_id: {burst_id}")
  if os.path.isfile(csvname_rtc_site_data):
    df_rtc_site = pd.read_csv(csvname_rtc_site_data)
  else:
    print("csv file not found - skipping plot")
    continue

  # Determine orbit direction
  df_burst_subset = df_burst.loc[df_burst['jpl_burst_id'] == burst_id]
  if len(df_burst_subset) == 1:
    odir = df_burst_subset['orbit_pass'].item()
  elif len(df_burst_subset) == 0:
    odir = 'NONE'
  else:
    odir = 'NONE'
    print('Warning: multiple burst id match')

  df_rtc_ts_wind = pd.read_csv(csvname_rtc_site_data).drop(columns=
    ['Unnamed: 0'])
  df_rtc_ts_wind['datetime'] = pd.to_datetime(df_rtc_ts_wind['datetime'])


  # Convert srting to list
  df_rtc_ts_wind['vv_l'] = df_rtc_ts_wind['vv'].apply(data_fcns.convert_to_list)
  df_rtc_ts_wind['vh_l'] = df_rtc_ts_wind['vh'].apply(data_fcns.convert_to_list)
  df_rtc_ts_wind['vv/vh_l'] = df_rtc_ts_wind['vv/vh'].apply(data_fcns.convert_to_list)

  df_rtc_ts_wind['vv_avg'] = df_rtc_ts_wind['vv_l'].apply(lambda x: sum(x) / len(x) if isinstance(x, list) else x)#.rolling(3).sum()
  df_rtc_ts_wind['vh_avg'] = df_rtc_ts_wind['vh_l'].apply(lambda x: sum(x) / len(x) if isinstance(x, list) else x)#.rolling(3).sum()
  df_rtc_ts_wind['vh_center'] = df_rtc_ts_wind['vh_l'].apply(lambda x: np.array(x[0:9]).reshape((3, 3))[1, 1] if isinstance(x, list) else x)#.rolling(3).sum()
  df_rtc_ts_wind['vv/vh_avg'] = df_rtc_ts_wind['vv/vh_l'].apply(lambda x: sum(x) / len(x) if isinstance(x, list) else x)

  figfile = 'tmp.png'
  figfile1 = 'tmp1.png'
  figfile2 = 'tmp2.png'
  try:
    plot_fcns.plot_rtc(df_rtc_ts_wind,figfile,df_site)
  except Exception as e:
    tb = sys.exc_info()[2]
    stack = inspect.trace()
    lcls = locals()
    slcls = tb.tb_next.tb_frame.f_locals
    print(f"Error in site id: {site_id}, burst_id: {burst_id}")

  try:
    #y_obs = alg_fcns.window_mahalanobis_1d_workflow1(df_rtc_ts_wind['vh_l'])
    vh_arrs = [np.array(w1).reshape((3,3)) for w1 in df_rtc_ts_wind['vh_l']]
    vv_arrs = [np.array(w1).reshape((3,3)) for w1 in df_rtc_ts_wind['vv_l']]
    vv_vh_arrs = [np.array(w1).reshape((3,3)) for w1 in df_rtc_ts_wind['vv/vh_l']]
    vv_vh_ts_obs,vv_vh_y_obs,vv_vh_dist1ds = alg_fcns.window_mahalanobis_1d_workflow2(vv_vh_arrs,acq_dt_l,threshold,n_pre_img,n_post_imgs_to_confirm,lookback_length_days)
    vv2_ts_obs,vv2_y_obs,vv2_dist1ds = alg_fcns.window_mahalanobis_2d_workflow2(vv_arrs,vh_arrs,acq_dt_l,threshold,n_pre_img,n_post_imgs_to_confirm,lookback_length_days)

    plt.rcParams['xtick.labelsize'] = 16

    fig, (ax1,ax2,ax3) = plt.subplots(3,1,sharex=True, figsize=(25,17))

    ax1.plot(df_rtc_ts_wind['datetime'], df_rtc_ts_wind['vv_avg'], marker='o', color='tab:blue', label='vv_avg')
    ax1.set_xlabel('Datetime')
    ax1.set_ylabel('vv_avg', color='tab:blue')
    ax1.tick_params(axis='y', labelcolor='tab:blue')
    ax1a = ax1.twinx()
    ax1a.plot(df_rtc_ts_wind['datetime'], df_rtc_ts_wind['vv/vh_avg'], marker='v', color='tab:brown', label='vv/vh_avg')
    ax1a.set_ylabel('vv/vh_avg', color='tab:brown')
    ax1a.tick_params(axis='y', labelcolor='tab:brown')
    change_type = df_site.change_type.iloc[0]
    plt.title(f'Change type {change_type}; {burst_id=}; {site_id=}; {odir=}',fontsize=20)
    ax1.grid(True)
    last_observed_time = df_site['last_observation_time'][0]
    if not isinstance(last_observed_time, type(pd.NaT)):
      ax1.axvline(x=last_observed_time, color='b', linestyle='--', label=f'Last observation time ({last_observed_time})')
    ax1a.legend(loc='upper left')
    change_time = df_site['change_time'][0]
    if not isinstance(change_time, type(pd.NaT)):
      ax1.axvline(x=change_time, color='r', linestyle='--', label=f'Change time ({change_time})')
    ax1.legend()

    ax2.scatter(vv_vh_ts_obs,vv_vh_y_obs,marker='s')
    ax2.set_xlabel('Acquisition Time of Post Image')
    ax2.set_ylabel('Change/No Change')
    ax2.set_ylim(-.05, 1.05)
    ax2r = ax2.twinx()
    ax2r.plot(vv_vh_ts_obs, vv_vh_dist1ds, marker='o', color='tab:blue', label='vh')
    ax2r.set_ylabel('Mahalanobis vv/vh', color='tab:brown')
    ax2r.tick_params(axis='y', labelcolor='tab:blue')
    ax2r.hlines(threshold,vv_vh_ts_obs.iloc[0],vv_vh_ts_obs.iloc[-1], ls='--')
    ax2.grid(True)

    #ax3.scatter(df_rtc_ts_wind['datetime'], y_obs)
    ax3.scatter(vv2_ts_obs,vv2_y_obs)
    ax3.tick_params(which='both', width=2)
    ax3.tick_params(which='major', length=7, labelsize=16)
    ax3.tick_params(which='minor', length=4, color='r')
    #ax3.set_xlim(df_rtc_ts_wind['datetime'].iloc[0],
    #  df_rtc_ts_wind['datetime'].iloc[-1])
    ax3.xaxis.set_major_locator(mdates.MonthLocator(interval=3))
    ax3.xaxis.set_major_formatter(mdates.ConciseDateFormatter(
      ax3.xaxis.get_major_locator()))
    ax3.xaxis.set_minor_locator(mdates.MonthLocator())
    ax3.set_ylabel('Change/No Change')
    ax3.set_ylim(-.05, 1.05)
    ax3r = ax3.twinx()
    bs = df_rtc_ts_wind['vh_center']#.rolling(3).sum()
    ax3r.plot(vv2_ts_obs, vv2_dist1ds, marker='v', color='tab:brown', label='vh')
    ax3r.set_ylabel('Mahalanobis 2D', color='tab:brown')
    ax3r.tick_params(axis='y', labelcolor='tab:brown')
    ax3r.hlines(threshold,vv2_ts_obs.iloc[0],vv2_ts_obs.iloc[-1], ls='--')
    ax3.grid(True)
    
    fig.tight_layout()
    fig.savefig(figfile2,dpi=300,bbox_inches="tight")
    plt.close(fig)

    
    slide = prs.slides.add_slide(blank_slide_layout_a)
    #if odir == 'ASCENDING':
    #  slide = prs_a.slides.add_slide(blank_slide_layout_a)
    #else:
    #  slide = prs_d.slides.add_slide(blank_slide_layout_d)

    left = Inches(0.5)
    top = Inches(1)
    #height = Inches(6)
    width = Inches(9)
    pic = slide.shapes.add_picture(figfile2,left,top,width=width)

  except Exception as e:
    exc_type,exc_value,exc_traceback = sys.exc_info()
    tb = exc_traceback
    stack = inspect.trace()
    lcls = locals()
    slcls = tb.tb_next.tb_frame.f_locals
    print(f"Alg Error in site id: {site_id}, burst_id: {burst_id}")
    print(exc_value)
    traceback.print_tb(tb)

prs.save('alg_cmp_2chan.pptx')
