#! /u/aurora-r0/richw/pkgs/miniforge3/envs/dist-s1/bin/python -i

from pptx import Presentation
from pptx.util import Inches
import pickle
import os
import sys
import inspect
import traceback
import numpy as np
import geopandas as gpd
import rasterio
from rasterio.crs import CRS
from rasterio.windows import Window
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import matplotlib
from shapely.geometry import Point
import pandas as pd
import geopandas as gpd
from dem_stitcher.geojson_io import read_geojson_gzip
from dem_stitcher.geojson_io import to_geojson_gzip
from shapely.geometry import Polygon, shape
import asf_search as asf
import concurrent.futures
from tqdm import tqdm
import itertools
from datetime import datetime, timedelta
from dateutil.parser import parse
import ast
import yaml
import einops

import distmetrics
import data_fcns
import plot_fcns
import test_setup
import alg_fcns
from alg_fcns import get_raster_fromfile
from alg_fcns import accuracy
from alg_fcns import anal_1d_alg
from alg_fcns import anal_2d_alg
from alg_fcns import anal_1d_logratio_aoi
from alg_fcns import anal_1d_mahalanobis_aoi
from alg_fcns import anal_2d_mahalanobis_aoi
from distmetrics import compute_mahalonobis_dist_1d
from distmetrics import compute_mahalonobis_dist_2d
from distmetrics import compute_transformer_zscore
from distmetrics import load_trained_transformer_model
from plot_fcns import prs_implot2
from plot_fcns import prs_roc4
from data_fcns import rasterize_shapes_to_array

# Speed up non-interactive plots
# Agg backend is optimized for non-interactive plots
matplotlib.use('Agg')
# Enable line simplification and increase the threshold
matplotlib.rcParams['path.simplify'] = True
matplotlib.rcParams['path.simplify_threshold'] = 1.0
# Split long lines to allow parallelization
matplotlib.rcParams['agg.path.chunksize'] = 10000

dir_base = '/home/richw/src/opera'
dir_research = dir_base + '/dist-s1-research'
dir_ad_hoc = dir_research + '/marshak/10_ad_hoc_data_generation'
dir_events_base = dir_base + '/dist-s1-events'
dir_events = dir_events_base + '/events'
dir_external_val_data = dir_events_base + '/external_validation_data_db'
#dir_events = dir_ad_hoc + '/events'
#dir_external_val_data = dir_ad_hoc + '/external_validation_data_db'
#dir_aurora_event_base = '/u/aurora-r0/cmarshak/dist-s1-research/marshak/10_ad_hoc_data_generation/out'
dir_aurora_event_base = '/u/aurora-r0/cmarshak/dist-s1-events/out'

event_name = 'chile_fire_2024'

yaml_file = dir_events + '/' + event_name + '.yml'
with open(yaml_file) as f:
  event_dict = yaml.safe_load(f)["event"]

print(f'Analyzing Event: {event_name}')

dir_event = dir_aurora_event_base + '/' + event_name
dir_rtc = dir_event + '/rtc_ts_merged'
dir_val = dir_event + '/validation_data'
dir_hls = dir_event + '/change_map_dist_hls'
dir_water = dir_event + '/water_mask'

filename_event = dir_external_val_data + '/' + event_name + '.geojson'
filename_event_db = event_name + '.png'
df_event = gpd.read_file(filename_event)
fig,ax = plt.subplots()
df_event.plot(ax=ax)
fig.tight_layout()
fig.savefig(filename_event_db,dpi=300,bbox_inches="tight")
plt.close(fig)

event_date = pd.to_datetime(event_dict['event_date'])

# Reference files available
ref_files = os.listdir(dir_val)
Nref = len(ref_files)
if Nref != 2:
  print("Warning: Not 2 files in validation directory")
for ival in range(Nref):
  filename_val = dir_val + '/' + ref_files[ival]
  #if ref_files[ival][0:6] == 'extent':
  if "extent" in ref_files[ival]:
    print("Loading AOI (extent) data")
    with rasterio.open(filename_val) as aoi:
      aoi_mask = aoi.read(1)
      aoi_mask = aoi_mask.astype(bool)
  else: 
    print("Loading Validation data")
    with rasterio.open(filename_val) as val:
      ref_profile = val.profile
      ref_crs = val.crs
      ref_mask = val.read(1)
      ref_mask = ref_mask.astype(bool)

df_event_utm = df_event.to_crs(ref_crs)

# ref and ref_mask are the same thing!
ref = rasterize_shapes_to_array(df_event_utm.geometry.tolist(), np.ones(df_event_utm.shape[0]), ref_profile, all_touched=True, dtype='uint8')

filename_ref = 'ref.png'
fig,ax = plt.subplots()
plt.imshow(ref, interpolation='none')
fig.tight_layout()
fig.savefig(filename_ref,dpi=300,bbox_inches="tight")
plt.close(fig)

filename_extent = 'extent.png'
fig,ax = plt.subplots()
plt.imshow(aoi_mask, interpolation='none')
fig.tight_layout()
fig.savefig(filename_extent,dpi=300,bbox_inches="tight")
plt.close(fig)

# DIST-HLS files available
print("Loading DIST-HLS data")
dist_hls_files = os.listdir(dir_hls)
Ndisthls = len(dist_hls_files)
dist_hls_data = [get_raster_fromfile(os.path.join(dir_hls,fname))
      for fname in dist_hls_files]

# Water msak files available
print("Loading water mask")
water_files = os.listdir(dir_water)
water_masks = [get_raster_fromfile(os.path.join(dir_water,fname))
      for fname in water_files]
land_mask = np.logical_not(water_masks[0])

# Alg parameters

# Tracks available for this event
tracknums = [d[5:] for d in os.listdir(dir_rtc)]

prs_dist1ds = Presentation()
#blank_slide_layout_dist1ds = prs_dist1ds.slide_layouts[6]
prs_data = Presentation()
#blank_slide_layout_data = prs_data.slide_layouts[6]
prs_dist2ds = Presentation()
prs_roc = Presentation()

algctrl = dict(
  td_lookback = timedelta(days=75),
  td_halfwindow = timedelta(days=50),
  Nconfirm = 1,
  do_data_plot = False,
  do_metric_plot = False,
  metric_plot_vmin = 0.0,
  metric_plot_vmax = 5.0,
  do_dist1ds_plot = False,
  do_dist2ds_plot = False,
  do_roc_plot = True,
  do_prc_plot = True,
  do_hist_mahalanobis_plot = False,
  plot_dir = 'plots')

file_mahalanobis_base = 'mahalanobis'

try:
  td_lookback = algctrl['td_lookback']
  td_halfwindow = algctrl['td_halfwindow']
  Nconfirm = algctrl['Nconfirm']
  do_data_plot = algctrl['do_data_plot']
  do_dist1ds_plot = algctrl['do_dist1ds_plot']
  do_dist2ds_plot = algctrl['do_dist2ds_plot']
  do_roc_plot = algctrl['do_roc_plot']
  plot_dir = algctrl['plot_dir']
  for tracknum in tracknums:
  #for tracknum in [tracknums[0]]:
    print(f"Track: {tracknum}")
    file_mahalanobis_track = file_mahalanobis_base + '_' + tracknum + '.pkl'
    dir_track = dir_rtc + '/track' + tracknum
    datetimes = []
    vv_filelist = []
    vh_filelist = []
    list_dir = sorted(os.listdir(dir_track))
    ipair = 0
    for filename_rtc in list_dir:
      # Accumulate datetimes and channel filenames
      # Assuming VV and VH pairs for each datetime
      str_channel = filename_rtc[-6:-4]
      str_datetime = filename_rtc[7:17]
      if str_channel == "VV":
        vv_filelist.append(filename_rtc)
      elif str_channel == "VH":
        vh_filelist.append(filename_rtc)
      else:
        print(f'Error: invalid str_channel {str_channel}')
        raise ValueError('Bad str_channel')
      if ipair == 0:
        # New datetime
        datetime1 = pd.to_datetime(str_datetime)
        datetimes.append(datetime1)
        ipair = 1
        last_str_datetime = str_datetime
      else:
        # Already have first member of pair
        ipair = 0
        if last_str_datetime != str_datetime:
          print(f'Error: mismatched datetime {str_datetime}')
          raise ValueError('Bad pair datetime')

    # Determine previous data list within available data times

    Ndatetimes = len(datetimes)
    pre_idxs = [alg_fcns.lookup_pre_idx_daterange(dt,datetimes,
      td_halfwindow,td_lookback) for i,dt in enumerate(datetimes)]

    # Set post indices to confirm

    post_idxs = [range(i,min(i+Nconfirm,Ndatetimes)) for i in range(Ndatetimes)]

    # Load all RTC data for available datetimes

    print("Loading RTC data")
    vv_data = [get_raster_fromfile(os.path.join(dir_track,fname))
      for fname in vv_filelist]
    vh_data = [get_raster_fromfile(os.path.join(dir_track,fname))
      for fname in vh_filelist]

    # Compute ratio
    vv_vh_ratio = [vv/vh for vv,vh in zip(vv_data,vh_data)]

    if do_data_plot:
      print("Data plot powerpoint")
      figfile = 'tmp.png'
      for i in range(Ndatetimes):
        print(f"{i+1}/{Ndatetimes}",end='\r')
        #if i > 0:
        #  continue
        prs_implot2('VV',vv_data[i],0.0,0.5,
          'VH',vh_data[i],0.0,0.05,
          datetimes[i],tracknum,event_dict['event_date'],
          prs_data,figfile)
      print("\ndone")

    print("Transformer")
    iuse = 8
    NN = 16
    model = load_trained_transformer_model()
    pre_vv = [vv_data[i][1000:1000+NN,1000:1000+NN] for i in pre_idxs[iuse]]
    pre_vh = [vh_data[i][1000:1000+NN,1000:1000+NN] for i in pre_idxs[iuse]]
    post_vv = vv_data[iuse][1000:1000+NN,1000:1000+NN]
    post_vh = vh_data[iuse][1000:1000+NN,1000:1000+NN]
    tmetric = compute_transformer_zscore(model,pre_vv,pre_vh,post_vv,post_vh,
      stride=2,batch_size=128)

    # Restrict to AOI for later accuracy calculations
    #ref_aoi = ref[aoi_mask]

    #thresholds = [x/5.0 for x in range(0,50)]
    #r0 = np.zeros_like(ref)
    #refs = [ref.view() if dt >= event_date else r0.view() for dt in datetimes]

    #dist1ds,change,tp,fp,tn,fn = anal_1d_alg(algctrl,datetimes,tracknum,
    #  event_date,
    #  'VV',vv_data,pre_idxs,post_idxs,
    #  thresholds,refs,land_mask)
    #dist1ds,change,tp,fp,tn,fn = anal_1d_alg_aoi2(algctrl,datetimes,tracknum,
    #  event_date,
    #  'VV',vv_data,pre_idxs,post_idxs,
    #  thresholds,ref_aoi,aoi_mask)

    print("Classical algorithms")
    thresholds = anal_1d_logratio_aoi(algctrl,datetimes,tracknum,
      event_date,'LR-VH',vh_data,pre_idxs,post_idxs,
      ref,aoi_mask)

    thresholds = anal_1d_logratio_aoi(algctrl,datetimes,tracknum,
      event_date,'LR-VV',vv_data,pre_idxs,post_idxs,
      ref,aoi_mask)

    thresholds = anal_1d_mahalanobis_aoi(algctrl,datetimes,tracknum,
      event_date,'VV',vv_data,pre_idxs,post_idxs,
      ref,aoi_mask)

    thresholds = anal_1d_mahalanobis_aoi(algctrl,datetimes,tracknum,
      event_date,'VH',vh_data,pre_idxs,post_idxs,
      ref,aoi_mask)

    thresholds = anal_2d_mahalanobis_aoi(algctrl,datetimes,tracknum,
      event_date,'2D_VV_VH',vv_data,vh_data,pre_idxs,post_idxs,
      ref,aoi_mask)

    #dist1ds,change,tp,fp,tn,fn = anal_2d_alg(algctrl,datetimes,tracknum,
    #  event_date,
    #  '2D_VV_VH',vv_data,vh_data,pre_idxs,post_idxs,
    #  thresholds,refs,land_mask)

except Exception as e:
    exc_type,exc_value,exc_traceback = sys.exc_info()
    tb = exc_traceback
    stack = inspect.trace()
    lcls = locals()
    tb1 = tb
    while True:
      # Sweep for bottom of traceback list within user code
      # (where the error occurred)
      if tb1.tb_next:
        fname = tb1.tb_next.tb_frame.f_code.co_filename
        if "opera" in fname:
          tb1 = tb1.tb_next
        else:
          break
      else:
        break
    sl = tb1.tb_frame.f_locals
    fname = tb1.tb_frame.f_code.co_filename
    lineno = tb1.tb_frame.f_lineno
    print(f"Error at tracknum: {tracknum}, filename_rtc: {filename_rtc}")
    print(f"Error at user line: {lineno} in {fname}")
    print(exc_value)
    #traceback.print_tb(tb)

if do_dist1ds_plot:
  prs_dist1ds.save('chile_fire_dist1ds.pptx')  

if do_dist2ds_plot:
  prs_dist2ds.save('chile_fire_dist2ds.pptx')  

if do_data_plot:
  prs_data.save('chile_fire_data.pptx')  

if do_roc_plot:
  prs_roc.save('chile_fire_roc.pptx')
