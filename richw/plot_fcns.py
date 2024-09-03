#! /usr/bin/env python

from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import pandas as pd
import data_fcns
import warnings
from datetime import datetime, timedelta
from dateutil.parser import parse

def plot_rtc(df_rtc_ts_wind,figfile,df_site):
  POL_RATIO_PLOT = False
  BURST_ID = df_rtc_ts_wind['burst_id'].iloc[0]
  SITE_ID = df_rtc_ts_wind['site_id'].iloc[0]
  fig, ax1 = plt.subplots(figsize=(10, 5))
  ax1.plot(df_rtc_ts_wind['datetime'], df_rtc_ts_wind['vv_avg'], marker='o', color='tab:blue', label='vv_avg')
 
  ax1.set_xlabel('Datetime')
  ax1.set_ylabel('vv_avg', color='tab:blue')
  ax1.tick_params(axis='y', labelcolor='tab:blue')
  ax2 = ax1.twinx()
  ax2.plot(df_rtc_ts_wind['datetime'], df_rtc_ts_wind['vh_avg'], marker='v', color='tab:brown', label='vh_avg')
  ax2.set_ylabel('vh_avg', color='tab:brown')
  ax2.tick_params(axis='y', labelcolor='tab:brown')
  if POL_RATIO_PLOT:
    ax3 = ax1.twinx()
    ax3.spines['right'].set_position(('outward', 60))  
    ax3.plot(df_rtc_ts_wind['datetime'], df_rtc_ts_wind['vv/vh_avg'], marker='P', color='tab:purple', label='vv/vh_avg')
    ax3.set_ylabel('vv/vh_avg', color='tab:purple')
    ax3.tick_params(axis='y', labelcolor='tab:purple')

  change_type = df_site.change_type.iloc[0]
  plt.title(f'Change type {change_type}; {BURST_ID=}; {SITE_ID=}')
  with warnings.catch_warnings():
    warnings.simplefilter("ignore", category=UserWarning)
    ax1.set_xticks(df_rtc_ts_wind['datetime'].tolist())
    ax1.set_xticklabels(df_rtc_ts_wind['datetime'].tolist(), rotation=90)

  # Subsample the xticks by N
  N = 2
  ticks = plt.gca().get_xticks()
  new_ticks = ticks[::N]
  new_tick_labels = [str(int(tick)) for tick in ticks[::N]]
  plt.gca().set_xticks(new_ticks)
  plt.gca().set_xticklabels(new_tick_labels)

  ax1.grid(True)

  last_observed_time = df_site['last_observation_time'][0]

  if not pd.isnull(last_observed_time):
    ax1.axvline(x=last_observed_time, color='b', linestyle='--', label=f'Last observation time ({last_observed_time})')
  ax2.legend(loc='upper left')

  change_time = df_site['change_time'][0]
  if not pd.isnull(change_time):
    ax1.axvline(x=change_time, color='r', linestyle='--', label=f'Change time ({change_time})')

  ax1.legend()
  fig.savefig(figfile,dpi=300,bbox_inches="tight")
  plt.close(fig)

def prs_implot2(chanstr1,arr1,vmin1,vmax1,chanstr2,arr2,vmin2,vmax2,
  dt,tracknum,event_date,prs,tmpname):
  fig,ax = plt.subplots()
  im = ax.imshow(arr1,cmap='gray',vmax=vmax1,vmin=vmin1)
  plt.title(f'{chanstr1}, trk: {tracknum}, {event_date}, im date: {datetime.strftime(dt,'%y-%m-%d')}')
  fig.tight_layout()
  fig.savefig(tmpname,dpi=300,bbox_inches="tight")
  plt.close(fig)
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  left = Inches(0.5)
  top = Inches(1)
  height = Inches(4)
  width = Inches(4)
  pic = slide.shapes.add_picture(tmpname,left,top,width=width,height=height)

  fig,ax = plt.subplots()
  im = ax.imshow(arr2,cmap='gray',vmax=vmax2,vmin=vmin2)
  plt.title(f'{chanstr2}, trk: {tracknum}, {event_date}, im date: {datetime.strftime(dt,'%y-%m-%d')}')
  fig.tight_layout()
  fig.savefig(tmpname,dpi=300,bbox_inches="tight")
  plt.close(fig)
  left = Inches(5)
  top = Inches(1)
  height = Inches(4)
  width = Inches(4)
  pic = slide.shapes.add_picture(tmpname,left,top,width=width,height=height)

def implot1(fig,ax,chanstr1,arr1,vmin1,vmax1,
  dt,tracknum,event_date,plotname):
  im = ax.imshow(arr1,cmap='gray',vmax=vmax1,vmin=vmin1)
  plt.title(f'{chanstr1}, trk: {tracknum}, {event_date}, im date: {datetime.strftime(dt,'%y-%m-%d')}')
  fig.tight_layout()
  fig.savefig(plotname,dpi=300,bbox_inches="tight")

def prs_roc4(chanstr1,tp1,fp1,chanstr2,tp2,fp2,chanstr3,tp3,fp3,
  chanstr4,tp4,fp4,dt,tracknum,event_date,prs,tmpname):
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  roc1(chanstr1,tp1,fp1,tmpname,dt,tracknum,event_date)
  left = Inches(0.5)
  top = Inches(1)
  height = Inches(3.5)
  width = Inches(4)
  pic = slide.shapes.add_picture(tmpname,left,top,width=width,height=height)
  roc1(chanstr2,tp2,fp2,tmpname,dt,tracknum,event_date)
  left = Inches(5)
  top = Inches(1)
  height = Inches(3.5)
  width = Inches(4)
  pic = slide.shapes.add_picture(tmpname,left,top,width=width,height=height)
  roc1(chanstr3,tp3,fp3,tmpname,dt,tracknum,event_date)
  left = Inches(0.5)
  top = Inches(4)
  height = Inches(3.5)
  width = Inches(4)
  pic = slide.shapes.add_picture(tmpname,left,top,width=width,height=height)
  roc1(chanstr4,tp4,fp4,tmpname,dt,tracknum,event_date)
  left = Inches(5)
  top = Inches(4)
  height = Inches(3.5)
  width = Inches(4)
  pic = slide.shapes.add_picture(tmpname,left,top,width=width,height=height)

def roc1(fig,ax,chanstr,tp1,fp1,tmpname,dt,tracknum,event_date): 
  line1 = ax.plot(fp1,tp1,marker='o')
  ax.set_xlabel('false positive')
  ax.set_ylabel('true positive')
  plt.title(f'{chanstr} trk: {tracknum}, {event_date}, {datetime.strftime(dt,'%y-%m-%d')}')
  fig.savefig(tmpname,dpi=300,bbox_inches="tight")
  
def hist1(fig,ax,chanstr,arr1,arr_name,nbins,binrange,
  tmpname,dt,tracknum,event_date): 
  ax.hist(arr1,bins=nbins,range=binrange)
  ax.set_xlabel(arr_name)
  ax.set_ylabel('count')
  plt.title(f'{chanstr} trk: {tracknum}, {event_date}, {datetime.strftime(dt,'%y-%m-%d')}')
  fig.savefig(tmpname,dpi=300,bbox_inches="tight")
  
def prs_dathist2(chanstr1,data1,name1,vmin1,vmax1,nbins1,binrange1,
  chanstr2,data2,name2,vmin2,vmax2,nbins2,binrange2,
  dt,tracknum,event_date,prs,tmpname):
  if binrange1 == 0:
    binrange1 = (data1.min(),data1.max())
  if binrange2 == 0:
    binrange2 = (data2.min(),data2.max())
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  fig,ax = plt.subplots()
  implot1(fig,ax,chanstr1,data1,vmin1,vmax1,dt,tracknum,event_date,tmpname)
  left = Inches(0.5)
  top = Inches(1)
  height = Inches(3.5)
  width = Inches(4)
  pic = slide.shapes.add_picture(tmpname,left,top,width=width,height=height)
  implot1(fig,ax,chanstr2,data2,vmin2,vmax2,dt,tracknum,event_date,tmpname)
  left = Inches(5)
  top = Inches(1)
  height = Inches(3.5)
  width = Inches(4)
  pic = slide.shapes.add_picture(tmpname,left,top,width=width,height=height)
  #plt.close(fig)

  #fig,ax = plt.subplots()
  ax.clear()
  hist1(fig,ax,chanstr1,data1,name1,nbins1,binrange1,
    tmpname,dt,tracknum,event_date)
  left = Inches(0.5)
  top = Inches(4)
  height = Inches(3.5)
  width = Inches(4)
  pic = slide.shapes.add_picture(tmpname,left,top,width=width,height=height)
  hist1(fig,ax,chanstr2,data2,name2,nbins2,binrange2,
    tmpname,dt,tracknum,event_date)
  left = Inches(5)
  top = Inches(4)
  height = Inches(3.5)
  width = Inches(4)
  pic = slide.shapes.add_picture(tmpname,left,top,width=width,height=height)
  plt.close(fig)

