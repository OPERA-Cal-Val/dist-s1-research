#! /u/aurora-r0/richw/pkgs/miniforge3/envs/dist-s1/bin/python -i

from pptx import Presentation
from pptx.util import Inches
import os
import sys
import inspect
import numpy as np
import geopandas as gpd
import rasterio
from rasterio.crs import CRS
from rasterio.windows import Window
from pathlib import Path
import matplotlib.pyplot as plt
from shapely.geometry import Point
import pandas as pd
from dem_stitcher.geojson_io import read_geojson_gzip
from dem_stitcher.geojson_io import to_geojson_gzip
from shapely.geometry import Polygon, shape
import asf_search as asf
import concurrent.futures
from tqdm import tqdm
import itertools
from datetime import datetime, timedelta
from dateutil.parser import parse
import ast
#import distmetrics
import data_fcns
import plot_fcns
import test_setup

dir_base = '/home/richw/src/opera'
try:
  df_rtc,df_val_bursts_subset,df_sites_subset,dir_rtc_site_data,dir_rtc_site_plots,basename_rtc_site_data = test_setup.setup(dir_base)
except:
  exc_type,exc_value,exc_traceback = sys.exc_info()
  tb = exc_traceback
  stack = inspect.trace()
  lcls = locals()
  slcls = tb.tb_next.tb_frame.f_locals
  raise Exception("Error in test_setup")


print('Iterating over site subset')

fstr = "%Y%m%dT%H%M%SZ"
fstr2 = "%Y-%m-%d %H:%M:%S"

delta_dt_all = []
delta_dt_all_maxabs = []
delta_vv_all = []
delta_vv_all_maxabs = []
delta_vh_all = []
delta_vh_all_maxabs = []
max_delta_dt_allsites = []
max_delta_vv_allsites = []
max_delta_vh_allsites = []
nan_count = []
lendat_mismatch_allsites = []
site_ids_compared = []
burst_ids_compared = []
delta_10min = datetime(2022,11,26,14,10,0) - datetime(2022,11,26,14,0,0)

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

Nsubset = len(df_val_bursts_subset)

#for index, df_row in df_val_bursts_subset.iterrows():
for index in range(Nsubset):
#for index in range(6,7):
#  index,df_row = next(df_val_bursts_subset.iterrows())
  # Need burst id using uppercase and dashes
  burst_id = df_val_bursts_subset.loc[index,'jpl_burst_id']
  site_id = df_val_bursts_subset.loc[index,'site_id']
  # if this csv file does not exist, it's data file will also have not data
  csvname_rtc_site_data = (dir_rtc_site_data + '/' + basename_rtc_site_data +
    str(site_id) + '_burst_' + burst_id + '.csv')
  plotname_rtc_site_data = (dir_rtc_site_plots + '/' + basename_rtc_site_data +
    str(site_id) + '_burst_' + burst_id + '.png')
  df_site = df_sites_subset[df_sites_subset.site_id
    == site_id].reset_index(drop=True)
  lon = df_site.geometry.x[0]
  lat = df_site.geometry.y[0]
  # Site locations use long,lat defined on WGS84 ellipsoid which is the basis
  # of the EPSG:4326 coordinate reference system (CRS)
  gdf1 = gpd.GeoDataFrame(df_site.geometry,crs="EPSG:4326")
  print(f"site id: {site_id}, burst_id: {burst_id}")
  if os.path.isfile(csvname_rtc_site_data):
    df_rtc_site = pd.read_csv(csvname_rtc_site_data)
  else:
    print("csv file not found - skipping plot")
    continue

  df_rtc_ts_wind = pd.read_csv(csvname_rtc_site_data).drop(columns=
    ['Unnamed: 0'])
  df_rtc_ts_wind['datetime'] = pd.to_datetime(df_rtc_ts_wind['datetime'])

  # Convert srting to list
  df_rtc_ts_wind['vv_l'] = df_rtc_ts_wind['vv'].apply(data_fcns.convert_to_list)
  df_rtc_ts_wind['vh_l'] = df_rtc_ts_wind['vh'].apply(data_fcns.convert_to_list)
  df_rtc_ts_wind['vv/vh_l'] = df_rtc_ts_wind['vv/vh'].apply(data_fcns.convert_to_list)

  df_rtc_ts_wind['vv_avg'] = df_rtc_ts_wind['vv_l'].apply(lambda x: sum(x) / len(x) if isinstance(x, list) else x)#.rolling(3).sum()
  df_rtc_ts_wind['vh_avg'] = df_rtc_ts_wind['vh_l'].apply(lambda x: sum(x) / len(x) if isinstance(x, list) else x)#.rolling(3).sum()
  df_rtc_ts_wind['vh_center'] = df_rtc_ts_wind['vh_l'].apply(lambda x: np.array(x[0:9]).reshape((3, 3))[1, 1] if isinstance(x, list) else x)#.rolling(3).sum()
  df_rtc_ts_wind['vv/vh_avg'] = df_rtc_ts_wind['vv/vh_l'].apply(lambda x: sum(x) / len(x) if isinstance(x, list) else x)

  figfile = 'tmp.png'
  try:
    plot_fcns.plot_rtc(df_rtc_ts_wind,figfile,df_site)
    slide = prs.slides.add_slide(blank_slide_layout)
    left = Inches(0.5)
    top = Inches(1)
    #height = Inches(6)
    width = Inches(9)
    #pic = slide.shapes.add_picture(figfile,left,top)
    pic = slide.shapes.add_picture(plotname_rtc_site_data,left,top,width=width)
  except Exception as e:
    tb = sys.exc_info()[2]
    stack = inspect.trace()
    lcls = locals()
    slcls = tb.tb_next.tb_frame.f_locals
    print(f"Error in site id: {site_id}, burst_id: {burst_id}")


prs.save('rtc.pptx')
